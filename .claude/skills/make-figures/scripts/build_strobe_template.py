#!/usr/bin/env python3
"""Build an editable STROBE participant flow diagram (.pptx) from a YAML config.

Why this script exists
----------------------
The Graphviz path (``generate_flow_diagram.R --type strobe``) renders an
auto-fitting monochrome diagram suitable for journal submission, but
co-authors often want to nudge box positions, edit prose, or recolor stage
labels in PowerPoint. This script produces a fully-editable .pptx in the
classical Identification → Screening → Inclusion → Analysis convention used
by the BMJ / Radiology / KJR cohort-study examples (Randolph 2018,
PLoS ONE 2021).

Unlike PRISMA, STROBE flow diagrams have study-specific spine structure
(number of cohort-stage boxes and exclusions varies per study), so this is
a single parametric builder rather than the PRISMA two-step
(``build_prisma2020_template.py`` then ``fill_prisma_template.py``).

YAML schema
-----------
    output_pptx: figures/figure1_strobe.pptx          # required
    slide_size: [10.5, 11.0]                           # inches; auto-computed if omitted
    title: "Figure 1. STROBE participant flow diagram"  # optional

    # `stages:` is OPTIONAL. STROBE itself does not mandate a phase column —
    # most cohort papers omit it (the column is a PRISMA 2020 convention). When
    # `stages:` is present, the script renders a left phase column and merges
    # consecutive same-stage spine rows under one label. When omitted (or empty),
    # the spine shifts left, the slide narrows, and the figure looks like the
    # plainer STROBE convention.
    stages:
      - {name: Identification, color: "#cfe1f5"}      # color = stage box fill
      - {name: Screening,      color: "#cfe1f5"}
      - {name: Inclusion,      color: "#cfe1f5"}
      - {name: Analysis,       color: "#cfe1f5"}
    spine:
      - {id: enrolled,  stage: Identification, text: "..."}
      - {id: screened,  stage: Screening,      text: "..."}
      - {id: eligible,  stage: Inclusion,      text: "..."}
      - {id: primary,   stage: Analysis,       text: "..."}
      - {id: landmark,  stage: Analysis,       text: "..."}  # consecutive same-stage rows share one stage label
    exclusions:
      - {after: enrolled, text: "Excluded (n = 1,200):\\n- did not meet the index finding on any screening"}
      - {after: screened, text: "Excluded (n = 3):\\n- prior outcome event / zero post-baseline follow-up"}

Usage
-----
    python3 build_strobe_template.py \
        --config figures/figure1_strobe.yaml \
        --out    figures/figure1_strobe.pptx

Open the resulting .pptx in PowerPoint to fine-tune positions or styling
before saving as PDF / TIFF for submission.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


# ── Defaults ────────────────────────────────────────────────────────────────
DEFAULT_SLIDE_W = 10.5    # default width when slide_size unset; height auto-computed
FALLBACK_SLIDE_H = 11.0   # only used if auto-compute can't fit
ABSOLUTE_MAX_SLIDE_H = 14.0  # PowerPoint hard upper bound for legible single-slide layouts

# Column geometry (inches) — tighter than widescreen; better for vertical flow
PHASE_X = 0.4
PHASE_W = 1.4
SPINE_X = 2.2
SPINE_W = 3.6
EXCL_X  = 6.4
EXCL_W  = 3.7

# Vertical spacing
TITLE_Y = 0.25
TITLE_H = 0.5
TOP_PAD = 0.35           # below title before first row
ROW_GAP = 0.35           # vertical gap between adjacent rows
DEFAULT_ROW_H = 1.20     # slightly taller boxes for better text breathing room
DEFAULT_EXCL_H = 1.05
BOTTOM_PAD = 0.4         # margin below last row

# Colors
NAVY   = RGBColor(0x1F, 0x3A, 0x68)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
LBLUE  = RGBColor(0xCF, 0xE1, 0xF5)


# ── Helpers ─────────────────────────────────────────────────────────────────
def _parse_color(hex_str: str | None, default: RGBColor) -> RGBColor:
    if not hex_str:
        return default
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _readable_text_color(bg: RGBColor) -> RGBColor:
    """Return BLACK or WHITE depending on background luminance."""
    # Relative luminance per WCAG (sRGB approx without gamma)
    r, g, b = bg[0], bg[1], bg[2]
    lum = 0.299 * r + 0.587 * g + 0.114 * b
    return RGBColor(0x1F, 0x3A, 0x68) if lum > 160 else WHITE


def add_box(slide, x, y, w, h, text, *,
            fill=WHITE, border=BLACK, font_color=BLACK,
            font_size=11, bold_first=True, anchor_middle=True,
            align_center=True, line_pt=1.0, dash=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(line_pt)
    if dash:
        ln = shape.line._get_or_add_ln()
        # Remove any existing prstDash to keep idempotent
        for existing in ln.findall(qn("a:prstDash")):
            ln.remove(existing)
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "dash")

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    if anchor_middle:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bool(bold_first and i == 0)
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=BLACK, width_pt=1.25):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    # Arrow head on the destination end
    line_xml = line.line._get_or_add_ln()
    tail = etree.SubElement(line_xml, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def add_right_angle_arrow(slide, x1, y1, x2, y2, *, color=BLACK, width_pt=1.0):
    """Spine right edge → exclusion left edge with a horizontal stub then arrow."""
    # Use a single straight horizontal connector since both ends share y
    add_arrow(slide, x1, y1, x2, y2, color=color, width_pt=width_pt)


# ── Loader ──────────────────────────────────────────────────────────────────
def load_config(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    if path.suffix in (".yaml", ".yml"):
        if not HAS_YAML:
            sys.exit("PyYAML not installed; install or use JSON config.")
        return yaml.safe_load(text)
    if path.suffix == ".json":
        return json.loads(text)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        if HAS_YAML:
            return yaml.safe_load(text)
        sys.exit("Config not JSON and PyYAML unavailable.")


# ── Builder ────────────────────────────────────────────────────────────────
def build(cfg: dict, out_path: Path) -> None:
    # Layout dimensions decided BEFORE slide creation so we can right-size the
    # canvas to the content (no large empty bottom strip).
    spine = cfg["spine"]
    n_spine = len(spine)
    spine_h = float(cfg.get("spine_box_height", DEFAULT_ROW_H))
    excl_h  = float(cfg.get("exclusion_box_height", DEFAULT_EXCL_H))
    title_text = cfg.get("title")

    # Phase column is optional. STROBE itself does not mandate one (PRISMA does);
    # most cohort papers omit it. Render only when the user supplies a non-empty
    # `stages:` list in the YAML config.
    stages_cfg = cfg.get("stages") or []
    has_phase_column = bool(stages_cfg)

    # Effective horizontal positions (collapse left margin when no phase column)
    if has_phase_column:
        spine_x_eff = SPINE_X
        excl_x_eff  = EXCL_X
    else:
        excl_gap = EXCL_X - SPINE_X - SPINE_W   # preserve the spine→excl gap
        spine_x_eff = PHASE_X
        excl_x_eff  = spine_x_eff + SPINE_W + excl_gap

    avail_top = TITLE_Y + TITLE_H + TOP_PAD if title_text else TITLE_Y
    row_pitch = spine_h + ROW_GAP
    last_row_end = avail_top + (n_spine - 1) * row_pitch + spine_h
    auto_h = last_row_end + BOTTOM_PAD
    auto_w = excl_x_eff + EXCL_W + PHASE_X      # right margin == left margin

    if "slide_size" in cfg:
        slide_w, slide_h = cfg["slide_size"]
    else:
        slide_w = max(auto_w, 4.0)
        slide_h = min(max(auto_h, 4.0), ABSOLUTE_MAX_SLIDE_H)

    # If the user gave a slide_size that's too short for the content, scale boxes down
    needed = last_row_end + BOTTOM_PAD
    if needed > slide_h:
        scale = (slide_h - avail_top - BOTTOM_PAD) / (n_spine * row_pitch - ROW_GAP)
        spine_h = max(0.7, spine_h * scale)
        excl_h  = max(0.6, excl_h * scale)
        row_pitch = spine_h + ROW_GAP

    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    if title_text:
        title_box = slide.shapes.add_textbox(
            Inches(PHASE_X), Inches(TITLE_Y),
            Inches(slide_w - PHASE_X * 2), Inches(TITLE_H),
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = NAVY

    spine_y = {}
    for i, b in enumerate(spine):
        spine_y[b["id"]] = avail_top + i * row_pitch

    # Phase column (optional) — group consecutive same-stage rows under one merged label.
    if has_phase_column:
        stage_color_lookup = {s["name"]: _parse_color(s.get("color"), NAVY) for s in stages_cfg}
        stages_seq = [b.get("stage") for b in spine]
        i = 0
        while i < n_spine:
            j = i
            while j + 1 < n_spine and stages_seq[j + 1] == stages_seq[i]:
                j += 1
            sname = stages_seq[i]
            y_top = spine_y[spine[i]["id"]]
            y_bot = spine_y[spine[j]["id"]] + spine_h
            h = y_bot - y_top
            stage_fill = stage_color_lookup.get(sname, NAVY)
            add_box(
                slide, PHASE_X, y_top, PHASE_W, h,
                sname,
                fill=stage_fill,
                border=stage_fill,
                font_color=_readable_text_color(stage_fill),
                font_size=14,
                bold_first=True,
                line_pt=0.0,
            )
            i = j + 1

    # Spine boxes + arrows
    for k, b in enumerate(spine):
        y = spine_y[b["id"]]
        add_box(
            slide, spine_x_eff, y, SPINE_W, spine_h,
            b["text"],
            fill=WHITE, border=BLACK, font_color=BLACK,
            font_size=11, bold_first=True,
        )
        if k > 0:
            prev_y = spine_y[spine[k - 1]["id"]]
            add_arrow(
                slide,
                spine_x_eff + SPINE_W / 2, prev_y + spine_h,
                spine_x_eff + SPINE_W / 2, y,
            )

    # Exclusion boxes + connector arrows
    # Vertically center each exclusion on its spine row's mid-height so the
    # connector arrow can be strictly horizontal (no diagonal segments).
    for excl in cfg.get("exclusions", []):
        after_id = excl["after"]
        y_after = spine_y[after_id]
        y_mid = y_after + spine_h / 2
        y_excl_top = y_mid - excl_h / 2
        add_box(
            slide, excl_x_eff, y_excl_top, EXCL_W, excl_h,
            excl["text"],
            fill=WHITE, border=BLACK, font_color=BLACK,
            font_size=10, bold_first=False,
            align_center=False,  # left-align text in exclusions
            dash=True,           # dashed border to visually distinguish exclusion side-branches
        )
        # Strictly horizontal connector: both endpoints at y_mid
        add_arrow(
            slide,
            spine_x_eff + SPINE_W, y_mid,
            excl_x_eff,            y_mid,
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


# ── CLI ─────────────────────────────────────────────────────────────────────
def main():
    p = argparse.ArgumentParser(description="Build an editable STROBE flow diagram .pptx from a YAML config.")
    p.add_argument("--config", required=True, type=Path, help="YAML/JSON STROBE config")
    p.add_argument("--out",    required=True, type=Path, help="Output .pptx path")
    args = p.parse_args()

    cfg = load_config(args.config)
    if not isinstance(cfg, dict):
        sys.exit(f"Config root must be a mapping; got {type(cfg)}")

    build(cfg, args.out)
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()
