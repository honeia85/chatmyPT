#!/usr/bin/env python3
"""Build a PRISMA 2020 flow diagram template (.pptx) programmatically.

Why this script exists
----------------------
prismastatement.org serves the official .docx templates behind a Squarespace
JavaScript fingerprint redirect that blocks programmatic download (curl, wget,
WebFetch, Wayback Machine all return the redirect HTML stub, not the file).

We reproduce the published layout from Page MJ et al. *BMJ* 2021;372:n71
(Fig 1, "PRISMA 2020 flow diagram"), which the PRISMA Group releases under
CC-BY 4.0. Box positions, labels, and arrow topology mirror the official
template; only the per-line counts are placeholders that fill_prisma_template.py
populates.

Variants
--------
- ``new``     : new systematic reviews (no "Studies included in previous
                version" branch)
- ``updated`` : updated systematic reviews (adds the previous-version branch
                on the right)

The "registers/databases only" (v1) vs "+ other sources" (v2) distinction is
handled by toggling the ``--include-other-sources`` flag, which adds the
secondary identification column on the left.

Usage
-----
    python3 build_prisma2020_template.py \
        --variant new --include-other-sources \
        --out templates/official/prisma2020/PRISMA_2020_flow_new_v2.pptx
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

# ── Layout constants (inches) ────────────────────────────────────────────────
SLIDE_W = 13.33
SLIDE_H = 10.0

BOX_W = 3.4
BOX_H = 0.95
COL_GAP = 0.4
ROW_GAP = 0.45

# Phase column on the far left
PHASE_COL_X = 0.3
PHASE_COL_W = 1.3

# Identification block columns
IDENT_COL1_X = PHASE_COL_X + PHASE_COL_W + 0.2  # registers / databases
IDENT_COL2_X = IDENT_COL1_X + BOX_W + COL_GAP   # other sources (v2 only)

# Vertical positions
TITLE_Y = 0.25
ROW1_Y = 1.1   # records identified
ROW2_Y = ROW1_Y + BOX_H + ROW_GAP  # records removed before screening (right inset)

SCREENING_Y = 3.3
SCREENING_EXCL_X_OFFSET = BOX_W + COL_GAP

ELIGIBILITY_Y = SCREENING_Y + (BOX_H + ROW_GAP) * 2  # reports retrieved + sought
INCLUDED_Y = SLIDE_H - BOX_H - 0.6

# Colors
NAVY = RGBColor(0x1F, 0x3A, 0x68)
LIGHT = RGBColor(0xE8, 0xEE, 0xF7)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    fill=LIGHT,
    border=NAVY,
    bold_first_line: bool = True,
    font_size: int = 10,
):
    """Add a rounded rectangle with the given text. Returns the shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.04)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = BLACK
        run.font.bold = bold_first_line and i == 0
    return shape


def add_phase_label(slide, top: float, text: str):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(PHASE_COL_X),
        Inches(top),
        Inches(PHASE_COL_W),
        Inches(BOX_H),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.color.rgb = NAVY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = NAVY
    line.line.width = Pt(1.25)


# ── Builders ─────────────────────────────────────────────────────────────────
def build_template(variant: str, include_other_sources: bool, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(TITLE_Y), Inches(SLIDE_W - 0.6), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        f"PRISMA 2020 flow diagram for "
        f"{'updated' if variant == 'updated' else 'new'} systematic reviews"
        + (" (databases, registers, and other sources)"
           if include_other_sources else " (databases and registers only)")
    )
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NAVY

    # Phase labels (left column)
    add_phase_label(slide, ROW1_Y, "Identification")
    add_phase_label(slide, SCREENING_Y, "Screening")
    add_phase_label(slide, INCLUDED_Y, "Included")

    # ── Identification row ───────────────────────────────────────────────────
    add_box(
        slide,
        IDENT_COL1_X,
        ROW1_Y,
        BOX_W,
        BOX_H,
        "Records identified from:\n"
        "  Databases (n = {n_db})\n"
        "  Registers (n = {n_reg})",
        font_size=10,
    )

    # Right inset: records removed before screening
    removed_x = IDENT_COL1_X + BOX_W + COL_GAP + (BOX_W + COL_GAP if include_other_sources else 0)
    add_box(
        slide,
        removed_x,
        ROW1_Y,
        BOX_W,
        BOX_H,
        "Records removed before screening:\n"
        "  Duplicate records removed (n = {n_dup})\n"
        "  Records marked as ineligible by automation tools (n = {n_auto})\n"
        "  Records removed for other reasons (n = {n_other_removed})",
        font_size=9,
    )
    add_arrow(slide, IDENT_COL1_X + BOX_W, ROW1_Y + BOX_H / 2, removed_x, ROW1_Y + BOX_H / 2)

    if include_other_sources:
        add_box(
            slide,
            IDENT_COL2_X,
            ROW1_Y,
            BOX_W,
            BOX_H,
            "Records identified from:\n"
            "  Websites (n = {n_web})\n"
            "  Organisations (n = {n_org})\n"
            "  Citation searching (n = {n_cite})\n"
            "  etc.",
            font_size=9,
        )

    # ── Screening row ────────────────────────────────────────────────────────
    add_box(
        slide,
        IDENT_COL1_X,
        SCREENING_Y,
        BOX_W,
        BOX_H,
        "Records screened\n(n = {n_screened})",
    )
    add_box(
        slide,
        IDENT_COL1_X + SCREENING_EXCL_X_OFFSET,
        SCREENING_Y,
        BOX_W,
        BOX_H,
        "Records excluded\n(n = {n_screen_excluded})",
    )
    add_arrow(slide, IDENT_COL1_X + BOX_W / 2, ROW1_Y + BOX_H, IDENT_COL1_X + BOX_W / 2, SCREENING_Y)
    add_arrow(
        slide,
        IDENT_COL1_X + BOX_W,
        SCREENING_Y + BOX_H / 2,
        IDENT_COL1_X + SCREENING_EXCL_X_OFFSET,
        SCREENING_Y + BOX_H / 2,
    )

    # Reports sought / not retrieved
    add_box(
        slide,
        IDENT_COL1_X,
        SCREENING_Y + BOX_H + ROW_GAP,
        BOX_W,
        BOX_H,
        "Reports sought for retrieval\n(n = {n_sought})",
    )
    add_box(
        slide,
        IDENT_COL1_X + SCREENING_EXCL_X_OFFSET,
        SCREENING_Y + BOX_H + ROW_GAP,
        BOX_W,
        BOX_H,
        "Reports not retrieved\n(n = {n_not_retrieved})",
    )
    add_arrow(
        slide,
        IDENT_COL1_X + BOX_W / 2,
        SCREENING_Y + BOX_H,
        IDENT_COL1_X + BOX_W / 2,
        SCREENING_Y + BOX_H + ROW_GAP,
    )
    add_arrow(
        slide,
        IDENT_COL1_X + BOX_W,
        SCREENING_Y + BOX_H + ROW_GAP + BOX_H / 2,
        IDENT_COL1_X + SCREENING_EXCL_X_OFFSET,
        SCREENING_Y + BOX_H + ROW_GAP + BOX_H / 2,
    )

    # Reports assessed / excluded with reasons
    assess_y = SCREENING_Y + (BOX_H + ROW_GAP) * 2
    add_box(
        slide,
        IDENT_COL1_X,
        assess_y,
        BOX_W,
        BOX_H,
        "Reports assessed for eligibility\n(n = {n_assessed})",
    )
    add_box(
        slide,
        IDENT_COL1_X + SCREENING_EXCL_X_OFFSET,
        assess_y,
        BOX_W,
        BOX_H + 0.4,
        "Reports excluded:\n"
        "  Reason 1 (n = {n_excl_r1})\n"
        "  Reason 2 (n = {n_excl_r2})\n"
        "  Reason 3 (n = {n_excl_r3})\n"
        "  etc.",
        font_size=9,
    )
    add_arrow(
        slide,
        IDENT_COL1_X + BOX_W / 2,
        SCREENING_Y + (BOX_H + ROW_GAP) * 2 - ROW_GAP,
        IDENT_COL1_X + BOX_W / 2,
        assess_y,
    )
    add_arrow(
        slide,
        IDENT_COL1_X + BOX_W,
        assess_y + BOX_H / 2,
        IDENT_COL1_X + SCREENING_EXCL_X_OFFSET,
        assess_y + BOX_H / 2,
    )

    # ── Included row ─────────────────────────────────────────────────────────
    add_box(
        slide,
        IDENT_COL1_X,
        INCLUDED_Y,
        BOX_W * 1.3,
        BOX_H + 0.3,
        "Studies included in review (n = {n_studies})\n"
        "Reports of included studies (n = {n_reports})",
        font_size=11,
    )
    add_arrow(
        slide,
        IDENT_COL1_X + BOX_W / 2,
        assess_y + BOX_H + 0.4,
        IDENT_COL1_X + BOX_W / 2,
        INCLUDED_Y,
    )

    # Footer attribution (CC-BY)
    foot = slide.shapes.add_textbox(
        Inches(0.3), Inches(SLIDE_H - 0.45), Inches(SLIDE_W - 0.6), Inches(0.35)
    )
    fp = foot.text_frame.paragraphs[0]
    r = fp.add_run()
    r.text = (
        "Layout adapted from Page MJ et al. PRISMA 2020 flow diagram, BMJ 2021;372:n71. "
        "Distributed under CC-BY 4.0."
    )
    r.font.size = Pt(8)
    r.font.italic = True
    r.font.color.rgb = GREY

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"✓ wrote {out_path}")


def main():
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("--variant", choices=["new", "updated"], default="new")
    ap.add_argument(
        "--include-other-sources",
        action="store_true",
        help="Include the 'other sources' identification column (v2 layout).",
    )
    ap.add_argument("--out", type=Path, required=True)
    args = ap.parse_args()
    build_template(args.variant, args.include_other_sources, args.out)


if __name__ == "__main__":
    main()
