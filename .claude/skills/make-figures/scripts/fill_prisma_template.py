#!/usr/bin/env python3
"""Populate count placeholders in a PRISMA 2020 .pptx template.

Reads a template built by ``build_prisma2020_template.py`` (or any .pptx whose
text frames contain ``{key}``-style tokens) and substitutes integer counts.

Two input modes
---------------
1. **Positional CSV** — ``--counts "315,122,186,7,111,204,102,84,3,15"``
   maps to a fixed sequence of keys appropriate for the most common DTA / IR
   meta-analysis workflow:

       n_db, n_dup, n_screened, n_screen_excluded,
       n_sought, n_assessed, n_excl_r1, n_excl_r2, n_excl_r3,
       n_studies

   This is a convenience for quick CLI use; remaining placeholders fall back
   to ``-`` so the figure remains visually balanced.

2. **JSON / YAML** — ``--counts-file counts.json`` provides a full mapping of
   every placeholder used in the template. Any keys not supplied render as
   ``-``.

Usage
-----
    python3 fill_prisma_template.py \
        --template templates/official/prisma2020/PRISMA_2020_flow_new_v1.pptx \
        --counts "315,122,186,7,111,204,102,84,3,15" \
        --out fig1_prisma_filled.pptx

    python3 fill_prisma_template.py \
        --template templates/official/prisma2020/PRISMA_2020_flow_new_v2.pptx \
        --counts-file my_counts.json \
        --out fig1_prisma_filled.pptx
"""
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from pptx import Presentation

POSITIONAL_KEYS = [
    "n_db",
    "n_dup",
    "n_screened",
    "n_screen_excluded",
    "n_sought",
    "n_assessed",
    "n_excl_r1",
    "n_excl_r2",
    "n_excl_r3",
    "n_studies",
]

ALL_KNOWN_KEYS = {
    "n_db", "n_reg", "n_dup", "n_auto", "n_other_removed",
    "n_web", "n_org", "n_cite",
    "n_screened", "n_screen_excluded",
    "n_sought", "n_not_retrieved",
    "n_assessed", "n_excl_r1", "n_excl_r2", "n_excl_r3",
    "n_studies", "n_reports",
}

TOKEN_RE = re.compile(r"\{([a-zA-Z_][a-zA-Z0-9_]*)\}")


def parse_counts(args) -> dict[str, str]:
    counts: dict[str, str] = {}
    if args.counts_file:
        data = json.loads(Path(args.counts_file).read_text())
        if not isinstance(data, dict):
            raise SystemExit("--counts-file must contain a JSON object mapping key→count")
        counts.update({k: str(v) for k, v in data.items()})
    if args.counts:
        values = [v.strip() for v in args.counts.split(",")]
        if len(values) != len(POSITIONAL_KEYS):
            raise SystemExit(
                f"--counts expects {len(POSITIONAL_KEYS)} comma-separated values "
                f"(got {len(values)}). Order: {', '.join(POSITIONAL_KEYS)}"
            )
        for k, v in zip(POSITIONAL_KEYS, values):
            counts.setdefault(k, v)
    return counts


def substitute(prs, counts: dict[str, str]) -> tuple[set[str], set[str]]:
    """Replace ``{key}`` tokens in every text run. Returns (filled, unfilled)."""
    filled: set[str] = set()
    unfilled: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    tokens = TOKEN_RE.findall(run.text)
                    if not tokens:
                        continue
                    new_text = run.text
                    for tok in tokens:
                        if tok in counts:
                            new_text = new_text.replace("{" + tok + "}", counts[tok])
                            filled.add(tok)
                        else:
                            new_text = new_text.replace("{" + tok + "}", "-")
                            unfilled.add(tok)
                    run.text = new_text
    return filled, unfilled


def main():
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("--template", type=Path, required=True)
    ap.add_argument("--counts", help="Comma-separated positional counts (10 values).")
    ap.add_argument("--counts-file", type=Path, help="JSON file mapping key→count.")
    ap.add_argument("--out", type=Path, required=True)
    args = ap.parse_args()

    if not args.counts and not args.counts_file:
        ap.error("supply --counts or --counts-file (or both)")

    counts = parse_counts(args)
    unknown = set(counts) - ALL_KNOWN_KEYS
    if unknown:
        print(f"⚠ unknown keys (will still substitute if matched): {sorted(unknown)}")

    prs = Presentation(args.template)
    filled, unfilled = substitute(prs, counts)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)

    print(f"✓ wrote {args.out}")
    print(f"  filled  : {sorted(filled)}")
    if unfilled:
        print(f"  blanks  : {sorted(unfilled)}  (rendered as '-')")


if __name__ == "__main__":
    main()
