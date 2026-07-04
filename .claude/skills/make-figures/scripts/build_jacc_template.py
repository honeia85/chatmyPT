"""Build the JACC Central Illustration PPTX template.

Layout matches official JACC submission templates (verified from
doi:10.1016/j.jacc.2019.10.035 Figures 1–4 PPTX files).

Run once. Output:
  references/visual_abstract_templates/jacc_central_illustration.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "references" / "visual_abstract_templates" / "jacc_central_illustration.pptx"

# JACC submission slide is 10 x 7.5 in (4:3 standard PPT)
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Blank layout
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# ---- Slot 1: citation text box ----
cite = slide.shapes.add_textbox(Inches(0.4), Inches(5.3), Inches(9.4), Inches(0.5))
tf = cite.text_frame
tf.text = "ARTICLECITATION — replace with: FirstAuthor et al. JACC: Asia YYYY; vol(issue):pages."
p = tf.paragraphs[0]
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

# ---- Slot 2: content picture placeholder (rectangle as marker) ----
ph = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(0.8), Inches(4.0), Inches(4.2)
)
ph.fill.solid()
ph.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
ph.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
ph.line.width = Pt(1.5)
ph.text_frame.text = "VISUALELEMENT — replace this rectangle with the author content figure (PNG/TIFF, 600 DPI, ~4×4.2 in)."
for para in ph.text_frame.paragraphs:
    para.font.size = Pt(10)
    para.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
    para.alignment = 1  # center

# ---- Slot 3: footer text box (reserved; usually empty in JACC templates) ----
foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(4.1), Inches(0.5))
foot.text_frame.text = "FOOTERNOTE — optional secondary caption (often empty)."
for para in foot.text_frame.paragraphs:
    para.font.size = Pt(9)
    para.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    para.font.italic = True

# ---- Slot 4: logo placeholder (rectangle as marker) ----
logo = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(6.7), Inches(2.7), Inches(0.8)
)
logo.fill.solid()
logo.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
logo.line.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
logo.line.width = Pt(0.75)
logo.text_frame.text = "JACCLOGO — JACC family logo (provided by editorial)."
for para in logo.text_frame.paragraphs:
    para.font.size = Pt(8)
    para.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    para.alignment = 1

# Save
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slide size: {Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in")
print("Placeholders: ARTICLECITATION, VISUALELEMENT, FOOTERNOTE, JACCLOGO")
