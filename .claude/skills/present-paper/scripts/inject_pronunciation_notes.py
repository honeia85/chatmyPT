#!/usr/bin/env python3
"""Append a pronunciation guide section to every slide's speaker notes.

Designed for non-native presenters who want a per-slide reading reference
without affecting the audience-facing view. Scans each slide's notes text
for tokens defined in a YAML/JSON ``PRON_DICT`` file (term → reading +
full-name), uses word-boundary regex to avoid false positives, and writes
a "[ 발음 ]" (or user-defined) section at the bottom of the notes.

Also auto-matches allele-style tokens (HLA-DRB1*07:01, HLA-A*02:01, …)
that match a configurable regex and synthesizes a reading by combining
the base reading from PRON_DICT with "스타 NN콜론NN".

This script is invocation-agnostic: it modifies a PPTX in place (or to a
new path) and does not assume any specific lecture topic. The PRON_DICT
file is supplied by the caller and is the only language/domain config.

Usage
-----

PRON_DICT is YAML or JSON. Each key is the term as it appears in the
notes text. Each value is a 2-tuple [reading, full_name].

```yaml
# pron_dict.yaml
GWAS: ["지와스", "Genome-wide association study"]
HLA: ["에이치-엘-에이", "Human leukocyte antigen"]
LGI1: ["엘-지-아이-원", "Leucine-rich glioma-inactivated 1"]
Perriot: ["페리오", "프랑스, t 묵음"]
```

```bash
python3 inject_pronunciation_notes.py input.pptx output.pptx \
    --dict pron_dict.yaml --header "[ 발음 ]"
```
"""
from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

DEFAULT_ALLELE_RE = r"\b(?:HLA-)?[A-Z]{1,5}[0-9]?\*[0-9]{2}:[0-9]{2}(?:N|L|S|Q)?\b"
DEFAULT_HEADER = "[ Pronunciation ]"


def load_dict(path: Path) -> dict:
    if path.suffix.lower() in (".yaml", ".yml"):
        try:
            import yaml  # type: ignore
        except ImportError:
            raise SystemExit("PyYAML required for YAML dict; pip install pyyaml")
        with path.open() as f:
            raw = yaml.safe_load(f)
    elif path.suffix.lower() == ".json":
        with path.open() as f:
            raw = json.load(f)
    else:
        raise SystemExit(f"unsupported dict format: {path.suffix}")
    out = {}
    for term, value in raw.items():
        if isinstance(value, (list, tuple)) and len(value) >= 1:
            reading = value[0]
            fullname = value[1] if len(value) > 1 else ""
        elif isinstance(value, str):
            reading, fullname = value, ""
        else:
            continue
        out[term] = (reading, fullname)
    return out


def find_terms(text: str, pron_dict: dict, allele_re: str | None):
    """Return [(term, reading, fullname)] for every dict key appearing in *text*."""
    if not text:
        return []
    hits = []
    seen = set()
    for term, (reading, fullname) in pron_dict.items():
        pat = r"(?<![A-Za-z0-9_])" + re.escape(term) + r"(?![A-Za-z0-9_])"
        if re.search(pat, text):
            if term not in seen:
                hits.append((term, reading, fullname))
                seen.add(term)
    if allele_re:
        allele_pattern = re.compile(allele_re)
        alleles = sorted(set(allele_pattern.findall(text)))
        for a in alleles:
            if a in seen:
                continue
            base = a.replace("HLA-", "").split("*")[0]
            tail = a.split("*", 1)[1] if "*" in a else ""
            base_reading = pron_dict.get(base, (base.lower(), ""))[0]
            reading = f"{base_reading} star {tail.replace(':', ' colon ')}"
            hits.append((a, reading, "HLA allele"))
            seen.add(a)
    return hits


def inject_notes(slide, terms, header: str):
    if not terms:
        return False
    tf = slide.notes_slide.notes_text_frame

    # blank separator
    p = tf.add_paragraph()
    r = p.add_run(); r.text = " "
    r.font.size = Pt(11)

    # header
    p = tf.add_paragraph()
    r = p.add_run(); r.text = header
    r.font.bold = True
    r.font.size = Pt(12)

    # each term: "▪ term  —  reading  ·  fullname"
    for term, reading, fullname in terms:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = f"▪ {term}"
        r.font.bold = True
        r.font.size = Pt(11)
        r2 = p.add_run(); r2.text = f"  —  {reading}"
        r2.font.size = Pt(11)
        if fullname:
            r3 = p.add_run(); r3.text = f"  ·  {fullname}"
            r3.font.italic = True
            r3.font.size = Pt(11)
    return True


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("src", type=Path)
    ap.add_argument("dst", type=Path)
    ap.add_argument("--dict", required=True, type=Path,
                    help="path to pron_dict.yaml or pron_dict.json")
    ap.add_argument("--header", default=DEFAULT_HEADER,
                    help="section header text (default: %(default)s)")
    ap.add_argument("--allele-regex", default=DEFAULT_ALLELE_RE,
                    help="regex for allele-style tokens (set to empty string to disable)")
    args = ap.parse_args()

    if not args.src.exists():
        print(f"source not found: {args.src}", file=sys.stderr)
        sys.exit(1)
    if args.src != args.dst:
        shutil.copy(args.src, args.dst)

    pron_dict = load_dict(args.dict)
    print(f"loaded {len(pron_dict)} terms from {args.dict}")

    allele_re = args.allele_regex if args.allele_regex else None
    prs = Presentation(args.dst)
    n_injected = 0
    n_terms_total = 0
    for slide in prs.slides:
        if not slide.has_notes_slide:
            continue
        body = slide.notes_slide.notes_text_frame.text
        if args.header in body:
            continue  # already injected on a previous run
        terms = find_terms(body, pron_dict, allele_re)
        if inject_notes(slide, terms, args.header):
            n_injected += 1
            n_terms_total += len(terms)
    prs.save(args.dst)
    print(f"injected pronunciation on {n_injected} slides ({n_terms_total} term entries)")
    print(f"OK: {args.dst}")


if __name__ == "__main__":
    main()
