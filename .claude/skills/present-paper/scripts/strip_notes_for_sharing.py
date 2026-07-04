#!/usr/bin/env python3
"""Strip all speaker notes from a PPTX file (sharing-ready variant).

Use case: after preparing an academic lecture deck where the speaker notes
contain presenter-only material (foreign-language narrative, pronunciation
hints, self-referential reminders), generate a clean copy whose notes are
empty so the deck can be circulated to the audience or to a senior reviewer
without leaking presenter-only content.

The slide body, figures, layout, and pronunciation/asterisk-bearing strings
(e.g. HLA alleles) are preserved exactly — only the notes_text_frame of
every slide is cleared. Also re-writes ``docProps/app.xml`` so the
PowerPoint Mac repair dialog is not triggered.

Usage:
    python3 strip_notes_for_sharing.py INPUT.pptx OUTPUT.pptx
"""
from __future__ import annotations

import argparse
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


def strip_notes(src: Path, dst: Path) -> dict:
    if src == dst:
        raise SystemExit("source and destination must differ")
    shutil.copy(src, dst)
    prs = Presentation(dst)
    cleared = 0
    for slide in prs.slides:
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf.text.strip():
                cleared += 1
            tf.clear()
    prs.save(dst)
    return {"cleared": cleared, "total_slides": len(prs.slides)}


def fix_app_xml(pptx_path: Path) -> dict:
    """Re-write docProps/app.xml so Slides/Notes counts match reality.

    python-pptx leaves Slides=0 on save, which PowerPoint Mac flags as a
    repair condition. Match canonical pattern in pptx-mac-compatibility.md §5.
    """
    prs = Presentation(pptx_path)
    n_slides = len(prs.slides)
    with zipfile.ZipFile(pptx_path, "r") as z:
        n_notes = sum(
            1 for x in z.namelist()
            if x.startswith("ppt/notesSlides/notesSlide") and x.endswith(".xml")
        )
    titles = []
    for slide in prs.slides:
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0][:60]
                break
        titles.append(title or "Untitled")

    def esc(t: str) -> str:
        return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    title_items = "".join(f"<vt:lpstr>{esc(t)}</vt:lpstr>" for t in titles)
    new_app_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        "<TotalTime>1</TotalTime><Words>0</Words>"
        "<Application>Microsoft Macintosh PowerPoint</Application>"
        "<PresentationFormat>Widescreen</PresentationFormat><Paragraphs>0</Paragraphs>"
        f"<Slides>{n_slides}</Slides><Notes>{n_notes}</Notes>"
        "<HiddenSlides>0</HiddenSlides><MMClips>0</MMClips><ScaleCrop>false</ScaleCrop>"
        '<HeadingPairs><vt:vector size="4" baseType="variant">'
        "<vt:variant><vt:lpstr>Theme</vt:lpstr></vt:variant>"
        "<vt:variant><vt:i4>1</vt:i4></vt:variant>"
        "<vt:variant><vt:lpstr>Slide Titles</vt:lpstr></vt:variant>"
        f"<vt:variant><vt:i4>{n_slides}</vt:i4></vt:variant>"
        "</vt:vector></HeadingPairs>"
        f'<TitlesOfParts><vt:vector size="{n_slides + 1}" baseType="lpstr">'
        f"<vt:lpstr>Office Theme</vt:lpstr>{title_items}"
        "</vt:vector></TitlesOfParts>"
        "<Manager></Manager><Company></Company><LinksUpToDate>false</LinksUpToDate>"
        "<SharedDoc>false</SharedDoc><HyperlinkBase></HyperlinkBase>"
        "<HyperlinksChanged>false</HyperlinksChanged><AppVersion>14.0000</AppVersion>"
        "</Properties>"
    )
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_path, "r") as zin, \
            zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            zout.writestr(item, new_app_xml if item == "docProps/app.xml" else zin.read(item))
    shutil.move(tmp, pptx_path)
    return {"slides": n_slides, "notes": n_notes}


def verify(pptx_path: Path) -> int:
    """Return the total number of non-whitespace characters left in any notes."""
    prs = Presentation(pptx_path)
    total = 0
    for slide in prs.slides:
        if slide.has_notes_slide:
            total += len(slide.notes_slide.notes_text_frame.text.strip())
    return total


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("src", type=Path)
    ap.add_argument("dst", type=Path)
    ap.add_argument("--no-app-xml-fix", action="store_true",
                    help="skip the app.xml rewrite (PowerPoint Mac may show repair dialog)")
    args = ap.parse_args()

    if not args.src.exists():
        print(f"source not found: {args.src}", file=sys.stderr)
        sys.exit(1)

    info = strip_notes(args.src, args.dst)
    print(f"cleared notes on {info['cleared']} / {info['total_slides']} slides")

    if not args.no_app_xml_fix:
        xml_info = fix_app_xml(args.dst)
        print(f"app.xml patched: Slides={xml_info['slides']}, Notes={xml_info['notes']}")

    leftover = verify(args.dst)
    if leftover > 0:
        print(f"WARNING: {leftover} chars of notes still remain", file=sys.stderr)
        sys.exit(2)
    print(f"OK: {args.dst}")


if __name__ == "__main__":
    main()
