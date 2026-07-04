#!/usr/bin/env python3
"""Inspect an institutional PowerPoint template (.pptx / .potx).

Lists every slide LAYOUT with its index, name, and placeholders (idx, type, name,
size), plus the theme fonts and first theme colors. Use the printed placeholder
indices to fill the template by `placeholder_format.idx` (SKILL.md "Mode C";
references/slide_visual_styles/institutional_brand.md) so the institution's master,
theme, and logo are preserved (patch-over-rebuild).

Usage:
    python3 inspect_pptx_template.py /path/to/template.potx
    python3 inspect_pptx_template.py /path/to/template.pptx --json

Read-only: never modifies the template.
"""
from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx required: pip install python-pptx")

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _emu_in(v) -> str:
    try:
        return f"{Emu(int(v)).inches:.2f}in"
    except Exception:
        return "?"


def theme_info(path: Path) -> dict:
    """Pull major/minor fonts + first scheme colors straight from theme1.xml."""
    out = {"fonts": {}, "colors": {}}
    try:
        with zipfile.ZipFile(path) as z:
            name = next((n for n in z.namelist() if n.startswith("ppt/theme/theme")), None)
            if not name:
                return out
            root = ET.fromstring(z.read(name))
        scheme = root.find(f"{A}themeElements/{A}fontScheme")
        if scheme is not None:
            for role in ("majorFont", "minorFont"):
                el = scheme.find(f"{A}{role}/{A}latin")
                if el is not None:
                    out["fonts"][role] = el.get("typeface", "")
        clr = root.find(f"{A}themeElements/{A}clrScheme")
        if clr is not None:
            for child in clr:
                tag = child.tag.split("}")[-1]
                srgb = child.find(f"{A}srgbClr")
                sys_ = child.find(f"{A}sysClr")
                val = (srgb.get("val") if srgb is not None
                       else (sys_.get("lastClr") if sys_ is not None else ""))
                if val:
                    out["colors"][tag] = f"#{val.upper()}"
    except Exception as e:  # pragma: no cover
        out["error"] = str(e)
    return out


def inspect(path: Path) -> dict:
    prs = Presentation(str(path))
    report = {
        "file": path.name,
        "slide_size": f"{prs.slide_width and Emu(prs.slide_width).inches:.2f} x "
                      f"{prs.slide_height and Emu(prs.slide_height).inches:.2f} in",
        "n_existing_slides": len(prs.slides),
        "theme": theme_info(path),
        "layouts": [],
    }
    for i, layout in enumerate(prs.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            pf = ph.placeholder_format
            phs.append({
                "idx": pf.idx,
                "type": str(pf.type).split(".")[-1].split(" ")[0] if pf.type is not None else "BODY",
                "name": ph.name,
                "size": f"{_emu_in(ph.width)} x {_emu_in(ph.height)}",
            })
        report["layouts"].append({"index": i, "name": layout.name, "placeholders": phs})
    return report


def render_text(r: dict) -> str:
    lines = [
        f"Template : {r['file']}",
        f"Size     : {r['slide_size']}    existing slides: {r['n_existing_slides']}",
    ]
    th = r.get("theme", {})
    if th.get("fonts"):
        lines.append(f"Fonts    : major={th['fonts'].get('majorFont','?')}  "
                     f"minor={th['fonts'].get('minorFont','?')}")
    if th.get("colors"):
        keys = ["dk1", "lt1", "dk2", "lt2", "accent1", "accent2"]
        shown = "  ".join(f"{k}={th['colors'][k]}" for k in keys if k in th["colors"])
        lines.append(f"Colors   : {shown}")
    lines.append("")
    lines.append("LAYOUTS (use index for add_slide, placeholder idx for filling):")
    for L in r["layouts"]:
        lines.append(f"  [{L['index']:2d}] {L['name']}")
        for ph in L["placeholders"]:
            lines.append(f"        idx={ph['idx']:<3} {ph['type']:<14} "
                         f"{ph['size']:<18} {ph['name']}")
    lines.append("")
    lines.append("Next: map outline slides to these layouts; fill by placeholder idx "
                 "(institutional_brand.md). Do NOT rebuild from scratch.")
    return "\n".join(lines)


def main() -> int:
    ap = argparse.ArgumentParser(description="Inspect a .pptx/.potx institutional template.")
    ap.add_argument("template", help="path to .pptx or .potx")
    ap.add_argument("--json", action="store_true", help="emit JSON instead of text")
    args = ap.parse_args()
    path = Path(args.template).expanduser()
    if not path.exists():
        sys.exit(f"not found: {path}")
    report = inspect(path)
    print(json.dumps(report, indent=2, ensure_ascii=False) if args.json else render_text(report))
    return 0


if __name__ == "__main__":
    sys.exit(main())
