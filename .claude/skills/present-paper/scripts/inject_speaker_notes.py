#!/usr/bin/env python3
"""Inject speaker notes into PowerPoint presentation slides.

This script adds or replaces speaker notes in a PPTX file without modifying
slide content, layout, or design. Notes are defined as a dictionary mapping
slide numbers (1-indexed) to note text.

Inline ``**bold**`` / ``*italic*`` in a note is parsed into run-level styling so the
asterisks do not show literally in Presenter View (python-pptx stores text verbatim).
Use --no-markdown for the legacy plain-text behavior.

Usage:
    python inject_speaker_notes.py input.pptx
    python inject_speaker_notes.py input.pptx -o output.pptx
    python inject_speaker_notes.py input.pptx --append
    python inject_speaker_notes.py input.pptx --dry-run
    python inject_speaker_notes.py input.pptx --no-markdown

Requirements:
    pip install python-pptx

License: MIT
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


# ---------------------------------------------------------------------------
# Speaker notes dictionary
# Map slide number (1-indexed) to note text.
# Empty string or missing key = skip that slide.
# ---------------------------------------------------------------------------
notes: dict[int, str] = {
    # 1: """Speaker note for slide 1.""",
    # 2: """Speaker note for slide 2.""",
}


# ---------------------------------------------------------------------------
# Inline-markdown rendering for notes.
# python-pptx stores ``notes_text_frame.text = "**bold**"`` VERBATIM, so the
# asterisks show literally in Presenter View. Parse ``**bold**`` / ``*italic*``
# (non-nested) into run-level styling instead. (Opt out with --no-markdown.)
# ---------------------------------------------------------------------------
_MD_INLINE = re.compile(r"(\*\*[^*\n]+\*\*|\*[^*\n]+\*)")


def _add_markdown_line(paragraph, line: str) -> None:
    """Emit one note line as styled runs, parsing **bold** / *italic*.

    Non-nested by design (matches the ``add_styled_note_line`` convention in
    ``pptx-speaker-notes.md``): malformed/mixed markers such as ``***x***`` render
    partially rather than nesting — text is never dropped and never raises.
    """
    if not line:
        paragraph.add_run().text = ""
        return
    for part in _MD_INLINE.split(line):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**") and len(part) > 4:
            run.text = part[2:-2]
            run.font.bold = True
        elif (part.startswith("*") and part.endswith("*")
              and not part.startswith("**") and len(part) > 2):
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part


def _render_notes_markdown(tf, text: str, append: bool) -> None:
    """Write text into the notes text frame with inline-markdown run styling.

    Preserves the line structure (one paragraph per line). With ``append`` and
    existing notes, inserts a ``---`` separator paragraph first.
    """
    lines = text.split("\n")
    if append and tf.text.strip():
        sep = tf.add_paragraph()
        sep.add_run().text = "---"
        for ln in lines:
            _add_markdown_line(tf.add_paragraph(), ln)
    else:
        tf.clear()  # leaves a single empty paragraph
        _add_markdown_line(tf.paragraphs[0], lines[0])
        for ln in lines[1:]:
            _add_markdown_line(tf.add_paragraph(), ln)


def inject_notes(
    input_path: str,
    output_path: str | None = None,
    append: bool = False,
    dry_run: bool = False,
    markdown: bool = True,
) -> None:
    """Inject speaker notes into a PPTX file.

    Args:
        input_path: Path to input PPTX file.
        output_path: Path to output PPTX file. Defaults to input with _notes suffix.
        append: If True, append to existing notes instead of replacing.
        dry_run: If True, print what would be done without saving.
    """
    input_file = Path(input_path)
    if not input_file.exists():
        print(f"Error: {input_file} not found")
        sys.exit(1)

    if output_path is None:
        output_file = input_file.with_stem(input_file.stem + "_notes")
    else:
        output_file = Path(output_path)

    prs = Presentation(str(input_file))
    total_slides = len(prs.slides)
    updated = 0

    for i, slide in enumerate(prs.slides, 1):
        if i not in notes or not notes[i]:
            continue

        if dry_run:
            preview = notes[i][:80].replace("\n", " ")
            print(f"  Slide {i:2d}: would {'append' if append else 'set'} → {preview}...")
            updated += 1
            continue

        if not slide.has_notes_slide:
            slide.notes_slide  # creates notes slide

        tf = slide.notes_slide.notes_text_frame
        if markdown:
            _render_notes_markdown(tf, notes[i], append)
        elif append and tf.text.strip():
            tf.text = tf.text + "\n\n---\n\n" + notes[i]
        else:
            tf.text = notes[i]
        updated += 1

    if dry_run:
        print(f"\nDry run: {updated}/{total_slides} slides would be updated")
        return

    prs.save(str(output_file))
    print(f"Done: {output_file} ({updated}/{total_slides} slides updated)")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Inject speaker notes into PowerPoint slides",
        epilog="Notes are defined in the 'notes' dictionary in this script.",
    )
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument(
        "-o", "--output",
        help="Output PPTX file (default: input with _notes suffix)",
    )
    parser.add_argument(
        "--append",
        action="store_true",
        help="Append to existing notes instead of replacing",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Print what would be done without saving",
    )
    parser.add_argument(
        "--no-markdown",
        action="store_true",
        help="Disable inline-markdown parsing (write **bold**/*italic* verbatim — legacy)",
    )
    args = parser.parse_args()

    if not notes:
        print("Warning: notes dictionary is empty. Edit this script to add notes.")
        print("Example:")
        print('  notes = {')
        print('      1: """Your note for slide 1.""",')
        print('      2: """Your note for slide 2.""",')
        print('  }')
        sys.exit(0)

    inject_notes(args.input, args.output, args.append, args.dry_run,
                 markdown=not args.no_markdown)


if __name__ == "__main__":
    main()
