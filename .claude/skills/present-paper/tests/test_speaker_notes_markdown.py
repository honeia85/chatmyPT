#!/usr/bin/env python3
"""Reproducible test for inject_speaker_notes.py inline-markdown rendering.

Builds a 1-slide PPTX, injects a note containing **bold** / *italic*, and asserts the
asterisks are gone and the runs are styled. --no-markdown must keep the text verbatim.
Skips cleanly (exit 0) when python-pptx is not installed. Network-free.

    python3 skills/present-paper/tests/test_speaker_notes_markdown.py
"""
import importlib.util
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed — SKIP (compile-only)")
    sys.exit(0)

ROOT = Path(__file__).resolve().parents[1]
SCRIPT = ROOT / "scripts" / "inject_speaker_notes.py"
spec = importlib.util.spec_from_file_location("isn", SCRIPT)
m = importlib.util.module_from_spec(spec)
spec.loader.exec_module(m)

tmp = Path(tempfile.mkdtemp())
prs = Presentation()
prs.slides.add_slide(prs.slide_layouts[6])
src = tmp / "in.pptx"
prs.save(str(src))

m.notes = {1: "Plain **bold** and *italic* here.\nSecond line."}
out = tmp / "in_notes.pptx"
m.inject_notes(str(src), str(out), markdown=True)
tf = Presentation(str(out)).slides[0].notes_slide.notes_text_frame
runs = [(r.text, r.font.bold, r.font.italic) for p in tf.paragraphs for r in p.runs]

out2 = tmp / "legacy.pptx"
m.inject_notes(str(src), str(out2), markdown=False)
tf2 = Presentation(str(out2)).slides[0].notes_slide.notes_text_frame

checks = [
    ("no literal '**' in rendered notes", "**" not in tf.text),
    ("no literal '*' in rendered notes", "*" not in tf.text),
    ("a bold run 'bold' exists", any(t == "bold" and b for t, b, i in runs)),
    ("an italic run 'italic' exists", any(t == "italic" and i for t, b, i in runs)),
    ("line structure preserved (>=2 paragraphs)", len(tf.paragraphs) >= 2),
    ("--no-markdown keeps '**bold**' verbatim", "**bold**" in tf2.text),
]
fail = 0
for name, ok in checks:
    print(("ok   " if ok else "FAIL ") + name)
    fail += 0 if ok else 1
print("ALL SPEAKER-NOTES MARKDOWN TESTS PASSED" if not fail else f"{fail} FAILED")
sys.exit(1 if fail else 0)
