#!/usr/bin/env python3
"""Nature/Lancet style PPTX builder — academic lecture / journal club / conference.

Generic template — pass slide content via builder functions. Speaker notes are
written in the user's preferred language (English by default; a Korean narrative
register is also supported for Korean presenters); slide body in English.

Color palette + typography + layout grid per:
  references/slide_visual_styles/nature_lancet.md

Mac PowerPoint compatibility per:
  ~/.claude/rules/pptx-mac-compatibility.md
  → run fix_app_xml() after final save (function included below)

Usage (illustrative):
    from build_pptx_nature_lancet import (
        new_presentation, add_title_slide, add_toc_slide, add_section_divider,
        add_transition_slide, add_content_slide, add_closing_slide,
        fix_app_xml,
    )

    prs = new_presentation()
    add_title_slide(prs,
        eyebrow="REVIEW LECTURE",
        title="<Lecture title in English>",
        subtitle="<One-line scope or lighthouse-paper subtitle>",
        meta_top="<Course>  ·  <Instructor>  ·  <Date>",
        meta_bottom="Presenter  <Name>  ·  <Department / Affiliation>",
        notes="<speaker notes in the user's preferred language>")

    add_toc_slide(prs, sections=[
        ("01", "<Section A title>", "<one-line summary>", "5 min"),
        ("02", "<Section B title>", "<one-line summary>", "8 min"),
        # ...
    ], notes="<...>", page_brand="<YEAR · COURSE NAME>")

    add_section_divider(prs, num="01", title="<Section A title>",
        subtitle="<lighthouse paper or section theme>", time_min=5)

    add_content_slide(prs,
        eyebrow="<COURSE / TOPIC LABEL>",
        title="<Sentence-headline title>",
        subtitle="<Author et al, Journal YYYY — n=...>",
        bullets=[
            "**Bold** main bullet",
            "  Sub bullet (prefix with 2 spaces)",
            "Another main bullet with *italic*",
        ],
        figure_path="figures/<your_fig>.png",
        fig_caption="<Short caption>",
        footnote="<Author et al, Journal YYYY>",
        page_brand="<YEAR · COURSE NAME>",
        notes="<notes in the user's preferred language>")

    add_closing_slide(prs, title="Take-home messages", bullets=[...], notes="...")

    out = Path("output.pptx")
    prs.save(out)
    fix_app_xml(out)
"""
from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree


# ============================================================================
# Style tokens — keep in sync with references/slide_visual_styles/nature_lancet.md
# ============================================================================

COLOR_NAVY        = RGBColor(0x1B, 0x2A, 0x4E)
COLOR_NAVY_LIGHT  = RGBColor(0x3F, 0x5A, 0x8C)
COLOR_TEXT        = RGBColor(0x21, 0x25, 0x2B)
COLOR_TEXT_SUB    = RGBColor(0x4A, 0x52, 0x5C)
COLOR_MUTED       = RGBColor(0x8A, 0x92, 0x9E)
COLOR_HAIRLINE    = RGBColor(0xCC, 0xD0, 0xD6)
COLOR_HIGHLIGHT   = RGBColor(0xB8, 0x3E, 0x3A)
COLOR_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG_SOFT     = RGBColor(0xF8, 0xF9, 0xFB)
COLOR_DIVIDER_BG  = RGBColor(0x12, 0x1D, 0x36)
COLOR_DIVIDER_SUB = RGBColor(0xC0, 0xCB, 0xDC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_KR = "Pretendard"
FONT_EN = "Inter"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ============================================================================
# Text helpers
# ============================================================================

def set_run(run, *, size: float = 20, bold: bool = False, italic: bool = False,
            color: RGBColor = COLOR_TEXT, ko: str = FONT_KR, en: str = FONT_EN,
            letter_space: str | None = None) -> None:
    """Apply Nature/Lancet font + size + color, with eastAsia attribute for Korean."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = en
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(f"{{{A_NS}}}ea"):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, f"{{{A_NS}}}ea")
    ea.set("typeface", ko)
    if letter_space is not None:
        rPr.set("spc", str(letter_space))


_MD_PATTERN = re.compile(r"(\*\*[^*\n]+\*\*|\*[^*\n]+\*)")


def add_styled(p, text: str, *, size: float = 20, color: RGBColor = COLOR_TEXT,
               ko: str = FONT_KR, en: str = FONT_EN) -> None:
    """Add a paragraph run with minimal markdown (**bold** / *italic*) parsed inline."""
    for part in _MD_PATTERN.split(text):
        if not part:
            continue
        r = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            r.text = part[2:-2]
            set_run(r, size=size, bold=True, color=color, ko=ko, en=en)
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            r.text = part[1:-1]
            set_run(r, size=size, italic=True, color=color, ko=ko, en=en)
        else:
            r.text = part
            set_run(r, size=size, color=color, ko=ko, en=en)


def add_notes(slide, text: str | None) -> None:
    """Inject speaker notes in the user's preferred language (with **bold**/*italic* parsed per run)."""
    if not text:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = ""
    for i, line in enumerate(text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if not line.strip():
            r = p.add_run(); r.text = " "
            set_run(r, size=12)
        else:
            add_styled(p, line, size=13)


# ============================================================================
# Presentation initialization
# ============================================================================

def new_presentation(width_in: float = 13.333, height_in: float = 7.5) -> Presentation:
    """Create a 16:9 widescreen Presentation with the Nature/Lancet defaults."""
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs


def _blank(prs: Presentation):
    """Return a blank-layout slide (layouts[6] is blank in default template)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================================
# Slide builders
# ============================================================================

def add_title_slide(prs: Presentation, *,
                    eyebrow: str,
                    title: str,
                    subtitle: str | None = None,
                    meta_top: str | None = None,
                    meta_bottom: str | None = None,
                    notes: str | None = None):
    s = _blank(prs)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.7), Inches(2.2), Emu(40000), Inches(3.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_NAVY
    bar.line.fill.background()

    eye = s.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(10), Inches(0.5))
    p = eye.text_frame.paragraphs[0]
    r = p.add_run(); r.text = eyebrow
    set_run(r, size=14, bold=True, color=COLOR_HIGHLIGHT, letter_space="300")

    box = s.shapes.add_textbox(Inches(1.1), Inches(2.7), Inches(11.5), Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    set_run(r, size=48, bold=True, color=COLOR_NAVY)
    p.space_after = Pt(8)
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        set_run(r2, size=18, italic=True, color=COLOR_TEXT_SUB)

    if meta_top or meta_bottom:
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.1), Inches(5.95), Inches(2.0), Emu(8000))
        line.fill.solid(); line.fill.fore_color.rgb = COLOR_NAVY
        line.line.fill.background()

        sub = s.shapes.add_textbox(Inches(1.1), Inches(6.0), Inches(11.5), Inches(1.0))
        stf = sub.text_frame; stf.word_wrap = True
        rows = [
            (meta_top,    16, True,  COLOR_NAVY),
            (meta_bottom, 13, False, COLOR_TEXT_SUB),
        ]
        first = True
        for text, sz, bld, col in rows:
            if not text:
                continue
            p = stf.paragraphs[0] if first else stf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(4)
            r = p.add_run(); r.text = text
            set_run(r, size=sz, bold=bld, color=col)

    add_notes(s, notes)
    return s


def add_section_divider(prs: Presentation, *,
                        num: str,
                        title: str,
                        subtitle: str | None = None,
                        time_min: int | None = None,
                        notes: str | None = None):
    s = _blank(prs)

    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_DIVIDER_BG
    bg.line.fill.background()

    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(1.2), Inches(3.0), Emu(15000), Inches(1.8))
    accent.fill.solid(); accent.fill.fore_color.rgb = COLOR_HIGHLIGHT
    accent.line.fill.background()

    num_box = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(3), Inches(1.5))
    p = num_box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = f"SECTION  {num}"
    set_run(r, size=18, bold=True, color=COLOR_HIGHLIGHT, letter_space="400")

    title_box = s.shapes.add_textbox(Inches(1.5), Inches(3.1), Inches(10), Inches(2.5))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    set_run(r, size=52, bold=True, color=COLOR_WHITE)
    p.space_after = Pt(8)
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        set_run(r2, size=20, italic=True, color=COLOR_DIVIDER_SUB)

    if time_min is not None:
        time_box = s.shapes.add_textbox(Inches(11.0), Inches(6.5), Inches(2.0), Inches(0.5))
        p = time_box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"{time_min}  MIN"
        set_run(r, size=13, bold=True, color=COLOR_MUTED, letter_space="500")

    add_notes(s, notes)
    return s


def add_transition_slide(prs: Presentation, *,
                          question: str,
                          notes: str | None = None):
    s = _blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_DIVIDER_BG
    bg.line.fill.background()

    box = s.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.8))
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER

    q = p.add_run(); q.text = "“ "
    set_run(q, size=72, bold=True, color=COLOR_HIGHLIGHT)
    r = p.add_run(); r.text = question
    set_run(r, size=36, bold=True, color=COLOR_WHITE)
    q2 = p.add_run(); q2.text = "  ”"
    set_run(q2, size=72, bold=True, color=COLOR_HIGHLIGHT)

    add_notes(s, notes)
    return s


def add_content_slide(prs: Presentation, *,
                       eyebrow: str,
                       title: str,
                       subtitle: str | None = None,
                       bullets: Sequence[str] = (),
                       figure_path: str | Path | None = None,
                       fig_caption: str | None = None,
                       footnote: str | None = None,
                       page_brand: str | None = None,
                       notes: str | None = None):
    """Standard content slide. Bullets prefixed with two spaces become sub-bullets."""
    s = _blank(prs)

    eye = s.shapes.add_textbox(Inches(0.7), Inches(0.32), Inches(8), Inches(0.4))
    p = eye.text_frame.paragraphs[0]
    r = p.add_run(); r.text = eyebrow
    set_run(r, size=10, bold=True, color=COLOR_MUTED, letter_space="300")

    title_box = s.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(12.0), Inches(1.1))
    ttf = title_box.text_frame; ttf.word_wrap = True
    p = ttf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    set_run(r, size=30, bold=True, color=COLOR_NAVY)
    p.space_after = Pt(3)
    if subtitle:
        p2 = ttf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle
        set_run(r2, size=15, italic=True, color=COLOR_TEXT_SUB)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(2.05), Inches(0.6), Emu(20000))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    line.line.fill.background()

    has_fig = figure_path is not None and Path(figure_path).exists()
    if has_fig:
        bullet_w = Inches(6.8)
        fig_x_in, fig_w_in, fig_h_in = 7.9, 5.0, 4.0
    else:
        bullet_w = Inches(12.0)

    bullet_box = s.shapes.add_textbox(Inches(0.7), Inches(2.4), bullet_w, Inches(4.6))
    btf = bullet_box.text_frame; btf.word_wrap = True
    for i, raw in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.line_spacing = 1.25

        is_sub = raw.startswith("  ")
        if is_sub:
            p.level = 1
            text = raw[2:]
            r = p.add_run(); r.text = "— "
            set_run(r, size=15, color=COLOR_MUTED)
            add_styled(p, text, size=16, color=COLOR_TEXT_SUB)
        else:
            r = p.add_run(); r.text = "▪  "
            set_run(r, size=14, bold=True, color=COLOR_HIGHLIGHT)
            add_styled(p, raw, size=20, color=COLOR_TEXT)

    if has_fig:
        from PIL import Image as PILImage
        img = PILImage.open(figure_path)
        iw, ih = img.size
        ar = iw / ih  # width / height
        if ar > fig_w_in / fig_h_in:
            w_in = fig_w_in
            h_in = fig_w_in / ar
        else:
            h_in = fig_h_in
            w_in = fig_h_in * ar

        fx_in = fig_x_in + (fig_w_in - w_in) / 2
        fy_in = 2.4 + (fig_h_in - h_in) / 2

        shadow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(fx_in + 0.06), Inches(fy_in + 0.06),
                                    Inches(w_in), Inches(h_in))
        shadow.fill.solid(); shadow.fill.fore_color.rgb = COLOR_HAIRLINE
        shadow.line.fill.background()

        pic = s.shapes.add_picture(str(figure_path),
                                   Inches(fx_in), Inches(fy_in),
                                   Inches(w_in), Inches(h_in))
        pic.line.color.rgb = COLOR_HAIRLINE
        pic.line.width = Emu(6000)

        if fig_caption:
            cap_y_in = min(fy_in + h_in + 0.10, 6.55)
            cap = s.shapes.add_textbox(Inches(fig_x_in), Inches(cap_y_in),
                                       Inches(fig_w_in), Inches(0.6))
            ctf = cap.text_frame; ctf.word_wrap = True
            cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run(); cr.text = "Figure  ·  "
            set_run(cr, size=10, bold=True, color=COLOR_HIGHLIGHT, letter_space="200")
            cr2 = cp.add_run(); cr2.text = fig_caption
            set_run(cr2, size=11, italic=True, color=COLOR_MUTED)

    if footnote:
        fn = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.35))
        fp = fn.text_frame.paragraphs[0]; fp.alignment = PP_ALIGN.RIGHT
        rr = fp.add_run(); rr.text = footnote
        set_run(rr, size=10, italic=True, color=COLOR_MUTED)

    if page_brand:
        pb = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(4.0), Inches(0.35))
        p = pb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = page_brand
        set_run(r, size=9, color=COLOR_MUTED, letter_space="300")

    add_notes(s, notes)
    return s


def add_toc_slide(prs: Presentation, *,
                   sections: Sequence[tuple[str, str, str, str]],
                   subtitle: str = "",
                   page_brand: str | None = None,
                   notes: str | None = None):
    """Outline slide. `sections` = [(num, title, summary, time_label), ...].

    Use num="·" for non-numbered final entry (e.g., wrap-up).
    """
    s = _blank(prs)

    eye = s.shapes.add_textbox(Inches(0.7), Inches(0.32), Inches(8), Inches(0.4))
    p = eye.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "OUTLINE"
    set_run(r, size=11, bold=True, color=COLOR_HIGHLIGHT, letter_space="400")

    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Outline"
    set_run(r, size=32, bold=True, color=COLOR_NAVY)
    if subtitle:
        p2 = tb.text_frame.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        set_run(r2, size=16, italic=True, color=COLOR_TEXT_SUB)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(2.1), Inches(0.6), Emu(20000))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    line.line.fill.background()

    y_start = 2.5
    row_h = 0.85
    for i, (num, sec_title, summary, time_label) in enumerate(sections):
        y = y_start + i * row_h

        nb = s.shapes.add_textbox(Inches(0.7), Inches(y), Inches(1.0), Inches(row_h))
        np_ = nb.text_frame.paragraphs[0]; np_.alignment = PP_ALIGN.LEFT
        nr = np_.add_run(); nr.text = num
        set_run(nr, size=22, bold=True,
                color=COLOR_HIGHLIGHT if num != "·" else COLOR_MUTED)

        tb2 = s.shapes.add_textbox(Inches(1.9), Inches(y), Inches(9.5), Inches(row_h))
        tt = tb2.text_frame; tt.word_wrap = True
        p = tt.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = sec_title
        set_run(r, size=22, bold=True, color=COLOR_NAVY)
        p.space_after = Pt(2)
        p2 = tt.add_paragraph()
        r2 = p2.add_run(); r2.text = summary
        set_run(r2, size=13, color=COLOR_TEXT_SUB)

        if time_label:
            time_b = s.shapes.add_textbox(Inches(11.5), Inches(y + 0.1),
                                          Inches(1.5), Inches(0.5))
            p = time_b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = time_label
            set_run(r, size=12, color=COLOR_MUTED)

        if i < len(sections) - 1:
            divider = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.7), Inches(y + row_h - 0.04),
                                         Inches(12.0), Emu(6000))
            divider.fill.solid(); divider.fill.fore_color.rgb = COLOR_HAIRLINE
            divider.line.fill.background()

    if page_brand:
        pf = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(12.0), Inches(0.35))
        pp_ = pf.text_frame.paragraphs[0]
        pr_ = pp_.add_run(); pr_.text = page_brand
        set_run(pr_, size=9, color=COLOR_MUTED, letter_space="300")

    add_notes(s, notes)
    return s


def add_glossary_slide(prs: Presentation, *,
                        eyebrow: str = "KEY ABBREVIATIONS",
                        title: str = "Glossary",
                        tier1: Sequence[tuple[str, str]] = (),
                        tier2: Sequence[tuple[str, str]] = (),
                        page_brand: str | None = None,
                        notes: str | None = None):
    """Glossary slide for multidisciplinary audiences.

    tier1: [(abbreviation, one-line context), ...]   — disease/concept (4–7 items)
    tier2: [(abbreviation, short definition), ...]   — method/statistics (8–12 items, rendered 2-column)

    See ~/.claude/rules/multidisciplinary-presentation.md §1.
    """
    s = _blank(prs)

    eye = s.shapes.add_textbox(Inches(0.7), Inches(0.32), Inches(8), Inches(0.4))
    p = eye.text_frame.paragraphs[0]
    r = p.add_run(); r.text = eyebrow
    set_run(r, size=10, bold=True, color=COLOR_MUTED, letter_space="300")

    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(12.0), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    set_run(r, size=30, bold=True, color=COLOR_NAVY)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.65), Inches(0.6), Emu(20000))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    line.line.fill.background()

    # Tier 1 — full width
    if tier1:
        box = s.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12.0), Inches(2.3))
        tf = box.text_frame; tf.word_wrap = True
        for i, (abbr, ctx) in enumerate(tier1):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
            r = p.add_run(); r.text = f"{abbr}  "
            set_run(r, size=15, bold=True, color=COLOR_NAVY)
            r2 = p.add_run(); r2.text = ctx
            set_run(r2, size=14, color=COLOR_TEXT_SUB)

    # Tier 2 — 2-column
    if tier2:
        half = (len(tier2) + 1) // 2
        for col, items in enumerate((tier2[:half], tier2[half:])):
            x_in = 0.7 + col * 6.3
            box = s.shapes.add_textbox(Inches(x_in), Inches(4.3),
                                        Inches(6.0), Inches(2.7))
            tf = box.text_frame; tf.word_wrap = True
            for i, (abbr, defn) in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(3)
                r = p.add_run(); r.text = f"{abbr}  "
                set_run(r, size=12, bold=True, color=COLOR_NAVY)
                r2 = p.add_run(); r2.text = defn
                set_run(r2, size=11, color=COLOR_TEXT_SUB)

    if page_brand:
        pb = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(4.0), Inches(0.35))
        p = pb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = page_brand
        set_run(r, size=9, color=COLOR_MUTED, letter_space="300")

    add_notes(s, notes)
    return s


def add_closing_slide(prs: Presentation, *,
                       title: str = "Take-home messages",
                       bullets: Sequence[str] = (),
                       contact: str | None = None,
                       page_brand: str | None = None,
                       notes: str | None = None):
    s = _blank(prs)

    eye = s.shapes.add_textbox(Inches(0.7), Inches(0.32), Inches(8), Inches(0.4))
    p = eye.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "CONCLUSION"
    set_run(r, size=11, bold=True, color=COLOR_HIGHLIGHT, letter_space="400")

    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(12.0), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    set_run(r, size=32, bold=True, color=COLOR_NAVY)

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.75), Inches(0.6), Emu(20000))
    line.fill.solid(); line.fill.fore_color.rgb = COLOR_HIGHLIGHT
    line.line.fill.background()

    if bullets:
        box = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.0))
        tf = box.text_frame; tf.word_wrap = True
        for i, raw in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(14)
            p.line_spacing = 1.3
            r = p.add_run(); r.text = "▪  "
            set_run(r, size=16, bold=True, color=COLOR_HIGHLIGHT)
            add_styled(p, raw, size=22, color=COLOR_TEXT)

    if contact:
        cb = s.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5))
        p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = contact
        set_run(r, size=12, color=COLOR_TEXT_SUB)

    if page_brand:
        pb = s.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(4.0), Inches(0.35))
        p = pb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = page_brand
        set_run(r, size=9, color=COLOR_MUTED, letter_space="300")

    add_notes(s, notes)
    return s


# ============================================================================
# Mac PowerPoint compatibility — app.xml fix
# ============================================================================

def fix_app_xml(pptx_path: str | Path) -> None:
    """Patch docProps/app.xml so Mac PowerPoint stops showing the repair dialog.

    Run after Presentation.save(). Updates <Slides>, <Notes>, HeadingPairs counts,
    and TitlesOfParts vector. See ~/.claude/rules/pptx-mac-compatibility.md §5.
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(pptx_path)
    n_slides = len(prs.slides)
    with zipfile.ZipFile(pptx_path, "r") as z:
        n_notes = sum(
            1 for x in z.namelist()
            if x.startswith("ppt/notesSlides/notesSlide") and x.endswith(".xml")
        )
    titles: list[str] = []
    for slide in prs.slides:
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().split("\n")[0][:60]
                break
        titles.append(title or "Untitled")

    esc = lambda t: t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    title_items = "".join(f"<vt:lpstr>{esc(t)}</vt:lpstr>" for t in titles)
    new_app_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" '
        'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        '<TotalTime>1</TotalTime><Words>0</Words>'
        '<Application>Microsoft Macintosh PowerPoint</Application>'
        '<PresentationFormat>Widescreen</PresentationFormat><Paragraphs>0</Paragraphs>'
        f'<Slides>{n_slides}</Slides><Notes>{n_notes}</Notes>'
        '<HiddenSlides>0</HiddenSlides><MMClips>0</MMClips><ScaleCrop>false</ScaleCrop>'
        '<HeadingPairs><vt:vector size="4" baseType="variant">'
        '<vt:variant><vt:lpstr>Theme</vt:lpstr></vt:variant>'
        '<vt:variant><vt:i4>1</vt:i4></vt:variant>'
        '<vt:variant><vt:lpstr>Slide Titles</vt:lpstr></vt:variant>'
        f'<vt:variant><vt:i4>{n_slides}</vt:i4></vt:variant>'
        '</vt:vector></HeadingPairs>'
        f'<TitlesOfParts><vt:vector size="{n_slides + 1}" baseType="lpstr">'
        f'<vt:lpstr>Office Theme</vt:lpstr>{title_items}'
        '</vt:vector></TitlesOfParts>'
        '<Manager></Manager><Company></Company><LinksUpToDate>false</LinksUpToDate>'
        '<SharedDoc>false</SharedDoc><HyperlinkBase></HyperlinkBase>'
        '<HyperlinksChanged>false</HyperlinksChanged><AppVersion>14.0000</AppVersion>'
        '</Properties>'
    )
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_path, "r") as zin, \
         zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            zout.writestr(item, new_app_xml if item == "docProps/app.xml" else zin.read(item))
    shutil.move(tmp, pptx_path)
    print(f"app.xml patched: Slides={n_slides}, Notes={n_notes}")


__all__ = [
    "new_presentation",
    "add_title_slide", "add_toc_slide", "add_section_divider",
    "add_transition_slide", "add_content_slide",
    "add_glossary_slide", "add_closing_slide",
    "set_run", "add_styled", "add_notes",
    "fix_app_xml",
    "COLOR_NAVY", "COLOR_HIGHLIGHT", "COLOR_TEXT", "COLOR_TEXT_SUB",
    "COLOR_MUTED", "COLOR_HAIRLINE", "COLOR_WHITE", "COLOR_DIVIDER_BG",
    "FONT_KR", "FONT_EN",
]
