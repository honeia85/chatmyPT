#!/usr/bin/env python3
"""Reference template library for editable academic PPTX generation.

Design tokens (NAVY/TEAL/ORANGE/GRAY, Apple SD Gothic Neo) and 16:9 fixed
coordinate zones are defaults — adapt to journal/institution conventions
as needed. NO markdown parsing, NO `cur_top` cumulative tracking (both
produce drift artifacts).

Each slide is defined as structured data → rendered by a type-specific
template function (T_lead / T_text / T_table / T_image_right / T_two_col
/ T_quote_slide / T_highlight_slide / T_metaphor_body / T_two_col_with_box
/ T_table_two_col).

`build_demo_slides()` at the bottom demonstrates every template with
generic academic placeholder content; copy and replace with your own
content function for a real deck.

Speaker-note injection: `parse_notes()` reads a markdown script with
`## Slide N:` headers and maps sections back to slide indices. For an
existing deck where you only want to add notes (no design regeneration),
use the sibling `inject_speaker_notes.py` instead.

Mac PowerPoint compatibility — verify before delivery:
  1. TIFF images: `find ppt/media -iname '*.tif*'` → convert to PNG
  2. `<a:sp3d>` 3D bevels in rPr render as red outlines on Mac → strip
  3. `docProps/app.xml` Slides/Notes/HeadingPairs/TitlesOfParts must
     match actual slide count
  4. `<a:srcRect>` values are 1/1000-percent (cap 100000) — never EMU

See `references/workflow-checklist.md` and the user-level rule
`~/.claude/rules/pptx-mac-compatibility.md` for the full checklist.
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Design tokens ──────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
TEAL    = RGBColor(0x00, 0x72, 0xB2)
ORANGE  = RGBColor(0xD5, 0x5E, 0x00)
GRAY    = RGBColor(0x33, 0x33, 0x33)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0x99, 0x99, 0x99)
BG_BLUE = RGBColor(0xF0, 0xF7, 0xFF)
BG_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
SUB_BLUE = RGBColor(0x7B, 0xB8, 0xE0)

FONT = 'Apple SD Gothic Neo'
SW = Inches(13.333)   # slide width
SH = Inches(7.5)      # slide height
ML = Inches(0.8)      # margin left
MR = Inches(0.8)      # margin right
MT = Inches(0.5)      # margin top
CW = SW - ML - MR     # content width

# Fixed vertical zones
TITLE_Y = MT
TITLE_H = Inches(0.8)
SUB_Y   = TITLE_Y + TITLE_H
SUB_H   = Inches(0.5)
BODY_Y  = SUB_Y + SUB_H + Inches(0.1)   # ~1.9in from top
BODY_H  = SH - BODY_Y - Inches(0.5)     # fill to bottom margin


# ── Low-level helpers ──────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _text(slide, x, y, w, h, txt, sz=20, color=GRAY, bold=False,
          align=PP_ALIGN.LEFT, auto_fit=False):
    """Add a text box. Returns the shape."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    _runs(p, txt, sz, color, bold)
    return tb

def _runs(para, txt, sz, color, bold):
    """Add text with **bold** inline markup support."""
    parts = re.split(r'(\*\*.*?\*\*)', txt)
    first = True
    for part in parts:
        if not part:
            continue
        is_b = part.startswith('**') and part.endswith('**')
        clean = part.strip('*') if is_b else part
        run = para.runs[0] if first and para.runs else para.add_run()
        first = False
        run.text = clean
        run.font.size = Pt(sz)
        run.font.color.rgb = NAVY if is_b else color
        run.font.bold = bold or is_b
        run.font.name = FONT

def _multiline(slide, x, y, w, h, lines, sz=20, color=GRAY, auto_fit=True):
    """Add multi-line text box with bullet support."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    first = True
    for line in lines:
        if not line.strip():
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)

        is_bullet = line.startswith('- ') or line.startswith('✓ ')
        txt = line.lstrip('-✓ ').strip() if is_bullet else line.strip()

        if line.strip().startswith('### '):
            _runs(p, line.strip()[4:], sz + 2, TEAL, True)
        else:
            if is_bullet:
                p.level = 0
            _runs(p, txt, sz, color, False)
    return tb

def _title_block(slide, title, subtitle=None, width=None):
    """Standard title + teal underline + optional subtitle."""
    w = width or CW
    _text(slide, ML, TITLE_Y, w, TITLE_H, title, sz=28, color=NAVY, bold=True)
    # Teal underline
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  ML, TITLE_Y + Inches(0.7), w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    if subtitle:
        _text(slide, ML, SUB_Y, w, SUB_H, subtitle, sz=22, color=TEAL)

def _table(slide, headers, rows, x, y, w, row_h=Inches(0.45)):
    """Add a styled table."""
    nr = len(rows) + 1
    nc = len(headers)
    if nc == 0:
        return
    shape = slide.shapes.add_table(nr, nc, x, y, w, row_h * nr)
    tbl = shape.table

    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h.strip('*')
        _cell_style(c, 15, WHITE, True, TEAL)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= nc:
                break
            c = tbl.cell(i + 1, j)
            c.text = val.strip('*')
            bg = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
            bld = val.strip().startswith('**')
            _cell_style(c, 14, GRAY, bld, bg)

def _cell_style(cell, sz, color, bold, bg):
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(sz)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tcPr = cell._tc.get_or_add_tcPr()
    sf = tcPr.makeelement(qn('a:solidFill'), {})
    sf.append(sf.makeelement(qn('a:srgbClr'), {'val': str(bg)}))
    tcPr.append(sf)

def _quote(slide, txt, x, y, w, h, sz=20):
    """Blockquote with teal left bar + blue background."""
    # Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.08), y, w - Inches(0.08), h)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_BLUE
    bg.line.fill.background()
    # Text
    _text(slide, x + Inches(0.3), y + Inches(0.1),
          w - Inches(0.5), h - Inches(0.2), txt, sz=sz, auto_fit=True)

def _highlight(slide, lines, x, y, w, h):
    """Yellow highlight box with orange border."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = BG_YELLOW
    box.line.color.rgb = ORANGE; box.line.width = Pt(2)
    _multiline(slide, x + Inches(0.3), y + Inches(0.2),
               w - Inches(0.6), h - Inches(0.4), lines, sz=18, auto_fit=True)

def _metaphor(slide, lines, x, y, w, h):
    """Light-blue metaphor box with teal left bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.08), y, w - Inches(0.08), h)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_BLUE
    bg.line.fill.background()
    _multiline(slide, x + Inches(0.3), y + Inches(0.15),
               w - Inches(0.6), h - Inches(0.3), lines, sz=17, auto_fit=True)

def _image(slide, path, x, y, w, h):
    """Add image preserving aspect ratio, centered in box."""
    if not os.path.exists(path):
        print(f'  WARNING: {path} not found')
        return
    from PIL import Image as PILImage
    with PILImage.open(path) as img:
        iw, ih = img.size
    aspect = iw / ih
    tw, th = int(w), int(h)
    if tw / th > aspect:
        fw, fh = int(th * aspect), th
    else:
        fw, fh = tw, int(tw / aspect)
    ax = x + (tw - fw) // 2
    ay = y + (th - fh) // 2
    slide.shapes.add_picture(path, ax, ay, fw, fh)

def _slidenum(slide, n, total):
    _text(slide, SW - Inches(1.2), SH - Inches(0.45), Inches(1.0), Inches(0.3),
          f'{n} / {total}', sz=10, color=LGRAY, align=PP_ALIGN.RIGHT)


# ── Slide type templates ───────────────────────────────────────
def T_lead(prs, title, subtitle='', extra=''):
    """Section divider: navy bg, centered white text."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _text(s, ML, Inches(2.2), CW, Inches(1.2), title,
          sz=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        _text(s, ML, Inches(3.6), CW, Inches(0.8), subtitle,
              sz=24, color=SUB_BLUE, align=PP_ALIGN.CENTER)
    if extra:
        _text(s, ML, Inches(4.6), CW, Inches(1.0), extra,
              sz=18, color=LGRAY, align=PP_ALIGN.CENTER)
    return s

def T_text(prs, title, body_lines, subtitle=None):
    """Title + bullet/text body. Most common type."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    _multiline(s, ML, BODY_Y, CW, BODY_H, body_lines, sz=20)
    return s

def T_table(prs, title, headers, rows, subtitle=None, body_before=None):
    """Title + optional body text + table."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    tbl_y = BODY_Y
    if body_before:
        _multiline(s, ML, BODY_Y, CW, Inches(1.5), body_before, sz=20)
        tbl_y = BODY_Y + Inches(1.6)
    _table(s, headers, rows, ML, tbl_y, CW)
    return s

def T_image_right(prs, title, body_lines, img_path, img_pct=30, subtitle=None):
    """Title + body on left, image on right."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    frac = img_pct / 100
    tw = int(SW * (1 - frac) - ML - Inches(0.3))
    iw = int(SW * frac)
    ix = SW - iw

    _title_block(s, title, subtitle, width=tw)
    _multiline(s, ML, BODY_Y, tw, BODY_H, body_lines, sz=20)
    _image(s, img_path, ix, Inches(0.2), iw, SH - Inches(0.4))
    return s

def T_quote_slide(prs, title, quotes, body_after=None, img_path=None,
                  img_pct=30, subtitle=None):
    """Title + blockquote(s). Optionally with image."""
    s = prs.slides.add_slide(prs.slide_layouts[6])

    if img_path:
        frac = img_pct / 100
        tw = int(SW * (1 - frac) - ML - Inches(0.3))
        iw = int(SW * frac)
        ix = SW - iw
        _title_block(s, title, subtitle, width=tw)
        _image(s, img_path, ix, Inches(0.2), iw, SH - Inches(0.4))
    else:
        tw = CW
        _title_block(s, title, subtitle)

    # Distribute quotes evenly in body area
    n = len(quotes)
    gap = Inches(0.15)
    q_h = int((BODY_H - gap * (n - 1)) / n) if n else BODY_H

    for i, q in enumerate(quotes):
        qy = BODY_Y + (q_h + gap) * i
        _quote(s, q, ML, qy, tw, q_h, sz=18)

    if body_after:
        _text(s, ML, SH - Inches(1.2), tw, Inches(0.6),
              body_after, sz=16, color=GRAY)
    return s

def T_two_col(prs, title, left_lines, right_lines, subtitle=None):
    """Two-column layout."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    cw = int((CW - Inches(0.5)) / 2)
    gap = Inches(0.5)
    _multiline(s, ML, BODY_Y, cw, BODY_H, left_lines, sz=18)
    _multiline(s, ML + cw + gap, BODY_Y, cw, BODY_H, right_lines, sz=18)
    return s

def T_two_col_with_box(prs, title, left_lines, right_lines, subtitle=None,
                       metaphor_col=None, metaphor_lines=None):
    """Two-column with metaphor/highlight box in one column."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    cw = int((CW - Inches(0.5)) / 2)
    gap = Inches(0.5)

    # Left column
    if metaphor_col == 'left' and metaphor_lines:
        _metaphor(s, metaphor_lines, ML, BODY_Y, cw, Inches(2.5))
        _multiline(s, ML, BODY_Y + Inches(2.7), cw, BODY_H - Inches(2.7),
                   left_lines, sz=18)
    else:
        _multiline(s, ML, BODY_Y, cw, BODY_H, left_lines, sz=18)

    # Right column
    rx = ML + cw + gap
    if metaphor_col == 'right' and metaphor_lines:
        _metaphor(s, metaphor_lines, rx, BODY_Y, cw, Inches(2.5))
        _multiline(s, rx, BODY_Y + Inches(2.7), cw, BODY_H - Inches(2.7),
                   right_lines, sz=18)
    else:
        _multiline(s, rx, BODY_Y, cw, BODY_H, right_lines, sz=18)
    return s

def T_highlight_slide(prs, title, highlight_lines, subtitle=None,
                      body_before=None):
    """Slide with prominent highlight box."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    hy = BODY_Y
    if body_before:
        _multiline(s, ML, BODY_Y, CW, Inches(1.0), body_before, sz=20)
        hy = BODY_Y + Inches(1.2)
    _highlight(s, highlight_lines, ML, hy, CW, SH - hy - Inches(0.5))
    return s

def T_metaphor_body(prs, title, body_lines, metaphor_lines, subtitle=None):
    """Body text + metaphor box below."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    mid = BODY_Y + BODY_H * 0.55
    _multiline(s, ML, BODY_Y, CW, mid - BODY_Y - Inches(0.1),
               body_lines, sz=20)
    _metaphor(s, metaphor_lines, ML, mid, CW, SH - mid - Inches(0.5))
    return s

def T_table_two_col(prs, title, left_lines, headers, rows, subtitle=None):
    """Left text + right table."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title_block(s, title, subtitle)
    cw = int((CW - Inches(0.5)) / 2)
    gap = Inches(0.5)
    _multiline(s, ML, BODY_Y, cw, BODY_H, left_lines, sz=18)
    _table(s, headers, rows, ML + cw + gap, BODY_Y, cw)
    return s


# ── Speaker notes ──────────────────────────────────────────────
def parse_notes(path):
    with open(path, 'r', encoding='utf-8') as f:
        content = f.read()
    # Match either `## Slide N:` (English) or legacy `--- 슬라이드 N: ... ---`
    pattern = r'(?:^##\s*Slide\s+(\d+)\s*:|---\s*슬라이드\s*(\d+)\s*:.*?---)'
    matches = list(re.finditer(pattern, content, flags=re.MULTILINE))
    notes = {}
    for i, m in enumerate(matches):
        num = int(m.group(1) or m.group(2))
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(content)
        raw = content[start:end].strip()
        # Clean
        lines = []
        for line in raw.split('\n'):
            s = line.strip()
            if s.startswith('## Part ') or s == '---' or s.startswith('```'):
                continue
            if s.startswith('*(') and s.endswith(')*'):
                lines.append(f'[{s[2:-2].strip()}]')
            else:
                lines.append(s.replace('**', '').replace('*', ''))
        text = re.sub(r'\n{3,}', '\n\n', '\n'.join(lines).strip())
        if text:
            notes[num] = text
    return notes


# ── Demo: showcase every template ───────────────────────────────
def build_demo_slides(prs):
    """Generic academic demo — one slide per template type.

    Replace this function with your own content builder. Keep all content
    as inline structured data (lists, dicts) — never parse markdown.
    """
    slides = []

    # 1. Title (T_lead)
    slides.append(T_lead(
        prs,
        'Editable Academic PPTX — Template Showcase',
        'A python-pptx reference deck',
        'Speaker · Affiliation · Date'))

    # 2. Section divider (T_lead)
    slides.append(T_lead(
        prs,
        'Section 1: Background',
        'Setting up the question'))

    # 3. Bullet body (T_text)
    slides.append(T_text(
        prs,
        'Background',
        ['- Disease X affects an estimated **N patients** worldwide.',
         '- Standard imaging modality: **modality Y**.',
         '- Existing limitations:',
         '  - Inter-reader variability',
         '  - Long acquisition time',
         '  - Reader fatigue on volume studies'],
        subtitle='Why this study matters'))

    # 4. Two-column compare (T_two_col)
    slides.append(T_two_col(
        prs,
        'Conventional vs. Proposed Approach',
        ['### Conventional',
         '- Manual segmentation',
         '- ~30 min per case',
         '- Reader-dependent'],
        ['### Proposed',
         '- Deep-learning assisted',
         '- ~3 min per case',
         '- Standardized output'],
        subtitle='Workflow comparison'))

    # 5. Quote slide (T_quote_slide)
    slides.append(T_quote_slide(
        prs,
        'Prior Evidence',
        ['"AI-assisted reading reduced inter-reader variability '
         'by 40% in a multi-center study." — Author et al., Journal Y, 2024',
         '"External validation remains the principal gap in clinical '
         'translation of imaging AI." — Reviewer A, 2025']))

    # 6. Table (T_table)
    slides.append(T_table(
        prs,
        'Cohort Characteristics',
        ['Variable', 'Training (n=800)', 'Internal test (n=200)', 'External (n=300)'],
        [['Age, mean (SD)', '62 (11)', '63 (10)', '60 (12)'],
         ['Female, n (%)', '420 (52.5)', '105 (52.5)', '160 (53.3)'],
         ['Disease prevalence', '38%', '40%', '35%'],
         ['**Imaging vendor mix**', 'A/B/C', 'A/B/C', 'D/E']],
        subtitle='Demographics across splits'))

    # 7. Image-right (T_image_right) — pass empty path to skip image safely
    slides.append(T_image_right(
        prs,
        'Methods Pipeline',
        ['- DICOM ingest with anonymization',
         '- Preprocessing: resampling, normalization',
         '- Model: 3D U-Net (or transformer variant)',
         '- Output: probability map + binary mask',
         '- QA: radiologist override loop'],
        img_path='figures/pipeline.png',  # replace with real asset
        img_pct=35,
        subtitle='From DICOM to decision'))

    # 8. Highlight slide (T_highlight_slide)
    slides.append(T_highlight_slide(
        prs,
        'Primary Result',
        ['**AUC = 0.92 (95% CI 0.89–0.95)** on external test set',
         'Sensitivity 0.88 · Specificity 0.91',
         'Non-inferior to senior reader (Δ AUC = 0.01, p = 0.78)'],
        subtitle='Diagnostic performance',
        body_before=['On the prospectively collected external cohort '
                     '(n = 300, 5 sites):']))

    # 9. Metaphor body (T_metaphor_body)
    slides.append(T_metaphor_body(
        prs,
        'Interpretation',
        ['- Performance preserved across vendor and protocol shifts.',
         '- Calibration check: ECE = 0.04 → clinically usable.',
         '- Failure modes concentrated in motion-degraded studies.'],
        ['Think of the model as a **second reader** — fast, consistent, '
         'never tired — but still requiring a final human sign-off on '
         'edge cases.'],
        subtitle='What it means in practice'))

    # 10. Two-col with metaphor box (T_two_col_with_box)
    slides.append(T_two_col_with_box(
        prs,
        'Limitations and Mitigations',
        ['### Limitations',
         '- Single-modality input',
         '- Retrospective external set',
         '- No outcome data yet'],
        ['### Planned next steps',
         '- Multi-modal fusion',
         '- Prospective registration',
         '- 12-month outcome follow-up'],
        metaphor_col='right',
        metaphor_lines=['**Open question** — is diagnostic accuracy a '
                        'sufficient endpoint, or do we need clinical '
                        'utility evidence?'],
        subtitle='Honest accounting'))

    # 11. Table-two-col (T_table_two_col)
    slides.append(T_table_two_col(
        prs,
        'Comparison with Published Models',
        ['### Take-aways',
         '- Comparable AUC at lower compute',
         '- Better calibration on minority subgroup',
         '- Released checkpoint + DICOM-SR exporter'],
        ['Model', 'AUC', 'Params'],
        [['Baseline', '0.85', '11M'],
         ['Prior SOTA', '0.91', '78M'],
         ['**Ours**', '**0.92**', '**24M**']],
        subtitle='Where this work sits'))

    # 12. Take-home (T_text)
    slides.append(T_text(
        prs,
        'Take-Home Messages',
        ['1. The model generalized to a **5-site external cohort** '
         '(AUC 0.92).',
         '2. Performance was **non-inferior to a senior reader**.',
         '3. **Calibration and failure-mode transparency** were the '
         'features that earned clinician trust — not raw AUC.',
         '4. Next: prospective validation tied to clinical outcomes.'],
        subtitle='Three points to remember'))

    return slides


# ── Main ────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    print('Building demo slides...')
    slides = build_demo_slides(prs)
    total = len(slides)
    print(f'  {total} slides created')

    # Slide numbers (skip lead/section-divider slides with solid bg)
    for i, slide in enumerate(slides):
        bg_type = slide.background.fill.type
        is_lead = bg_type is not None and 'SOLID' in str(bg_type)
        if not is_lead:
            _slidenum(slide, i + 1, total)

    # Optional speaker-note injection — uncomment if a script file exists
    # notes = parse_notes('speaker_script.md')
    # for i, slide in enumerate(slides):
    #     num = i + 1
    #     if num in notes:
    #         slide.notes_slide.notes_text_frame.text = notes[num]

    out = 'demo_slide_deck.pptx'
    prs.save(out)
    sz = os.path.getsize(out) / 1024
    print(f'\nDone: {out}')
    print(f'  Slides: {total}')
    print(f'  Size: {sz:.0f} KB, Editable: YES')


if __name__ == '__main__':
    main()
